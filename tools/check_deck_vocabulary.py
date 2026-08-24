#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Catch names WE introduced that the source deck never uses.

THE FAILURE THIS EXISTS FOR. CMS Lecture 9 slide 50 has no text at all -- it is
purely a figure of five skin blocks labelled "Level I" to "Level V". The word
"Clark" appears NOWHERE in that deck. I nonetheless wrote "Clark level" through
the guide, the cram sheet, four quiz questions and five Arcade cards, as though
it were the deck's own wording.

The content was correct; the sourcing was not. Two harms, both real for a
student:

  - Searching their slides for the term finds nothing, and the reasonable
    conclusion is that they have lost a slide.
  - An exam written from the deck will use the DECK's word. Studying a synonym
    is fine; not knowing it is a synonym is not.

This checks a watchlist of terms that are easy to import by reflex from general
medical knowledge. For each, if the term appears in our content for a class but
NOT in that class's decks, it must be introduced explicitly -- the content has
to say somewhere that the deck does not use it.

Run with a class directory, e.g.
  python3 tools/check_deck_vocabulary.py "Clinical Medicine and Surgery I Exam 1" \\
      ~/Desktop/"Semester 2"/"Clinical Medicine and Surgery I Inbox"/"Exam 1"
"""
import os, re, sys, glob

# Terms a writer reaches for automatically. Extend as new ones are caught.
WATCH = ["Clark", "Breslow", "TNM", "Fitzpatrick", "Nikolsky", "Auspitz",
         "Koebner", "Wickham", "Gottron", "Darier", "Asboe-Hansen",
         "Tzanck", "Wood lamp", "ABCDE", "Hutchinson"]

# Phrasings that count as introducing a term rather than assuming it.
INTRODUCED = re.compile(
    r"does not appear anywhere in this deck|the deck never says|"
    r"conventionally called|the conventional name|deck labels these|"
    r"deck&rsquo;s own wording|not the deck's term|"
    # A DISTRACTOR that names something absent, with an explanation saying so,
    # is a legitimate use -- arguably the best one, since it teaches that the
    # thing is not on the slide. Two real cases were flagged: "Wood lamp is not
    # on this slide's bedside list" and "the stated reason is Breslow depth,
    # not Clark level".
    r"is not on this slide|not on this slide's|is not on the list|"
    r"the stated reason is|are not on this list|is not one of the", re.I)


def deck_text(inbox):
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        sys.exit("python-pptx is required -- install it rather than skipping this check")
    def walk(shapes):
        for sh in shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from walk(sh.shapes)
            else:
                yield sh
    out = []
    decks = glob.glob(os.path.join(inbox, "*.pptx"))
    assert decks, "no .pptx found in %s -- a silent empty scan would pass everything" % inbox
    for d in decks:
        p = Presentation(d)
        for s in p.slides:
            for sh in walk(s.shapes):
                if sh.has_text_frame:
                    out.append(sh.text_frame.text)
            if s.has_notes_slide:
                out.append(s.notes_slide.notes_text_frame.text)
    return " ".join(out).lower(), len(decks)


def main():
    if len(sys.argv) < 3:
        sys.exit(__doc__)
    classdir, inbox = sys.argv[1], os.path.expanduser(sys.argv[2])
    decks, ndecks = deck_text(inbox)

    ours = []
    for f in glob.glob(os.path.join(classdir, "*.html")):
        ours.append((os.path.basename(f), open(f, encoding="utf-8").read()))
    assert ours, "no HTML found in %s" % classdir

    problems = []
    for term in WATCH:
        t = term.lower()
        in_deck = t in decks
        for name, body in ours:
            if t not in body.lower():
                continue
            if in_deck:
                continue
            # THE EXEMPTION MUST BE SCOPED TO THE TERM, NOT THE PAGE. Checking
            # INTRODUCED against the whole page meant one caveat anywhere on it
            # excused every watchlist term on that page -- so a correctly
            # introduced "Clark" would have silently covered for an
            # unintroduced second term. The caveat now has to sit near an
            # occurrence of the term it is excusing.
            low = body.lower()
            excused = False
            for m in re.finditer(re.escape(t), low):
                window = body[max(0, m.start() - 600): m.start() + 600]
                if INTRODUCED.search(window):
                    excused = True
                    break
            if excused:
                continue
            problems.append((term, name))

    print("scanned %d deck(s) and %d page(s) for %d watchlist term(s)"
          % (ndecks, len(ours), len(WATCH)))
    if problems:
        print("\nTERMS USED IN OUR CONTENT BUT ABSENT FROM THE DECKS, and not introduced as such:")
        for term, name in problems:
            print("   %-14s in %s" % (term, name))
        print("\nEither use the deck's own wording, or say plainly that the deck does not use this term.")
        sys.exit(1)
    print("every watchlist term either appears in the decks or is explicitly introduced")


if __name__ == "__main__":
    main()
